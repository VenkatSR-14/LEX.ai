from pptx import Presentation
import openai
import os

def extract_text_from_file(file_path):
    """
    Extract text from a file.
    For PPTX files, use python-pptx; for other files, attempt to read as plain text.
    """
    if file_path.lower().endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            return ""

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        return ""

def segment_text(text):
    """
    Segment text into chapters/sections using a regex.
    Chapters are defined by lines that start with "Chapter" or "Section".
    """
    import re
    chapters = {}
    current_chapter = "Introduction"
    chapters[current_chapter] = []
    lines = text.splitlines()
    pattern = re.compile(r'^(Chapter|Section)\s+\S+', re.IGNORECASE)
    for line in lines:
        if pattern.match(line):
            current_chapter = line.strip()
            chapters[current_chapter] = []
        else:
            chapters[current_chapter].append(line)
    for key in chapters:
        chapters[key] = "\n".join(chapters[key]).strip()
    return chapters



import openai

def generate_explanation(chapter_title, content):
    """
    Generate a detailed explanation for a chapter using OpenAI's GPT-3.5-turbo.
    Ensure that your environment variable OPENAI_API_KEY is set.
    """
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that explains content in detail."},
                {"role": "user", "content": f"Generate a detailed explanation for the following chapter titled '{chapter_title}' based on the content below:\n\n{content}"}
            ],
            temperature=0.7,
            max_tokens=300
        )
        explanation = response.choices[0].message["content"].strip()
        return explanation
    except Exception as e:
        return f"Error generating explanation: \n\n{str(e)}"
